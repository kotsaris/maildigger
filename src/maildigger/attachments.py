"""Attachment saving and format conversion to LLM-readable text."""

import csv
import io
from pathlib import Path

from rich.console import Console

from .parse import AttachmentInfo

console = Console(stderr=True)


def save_and_convert(
    attachment: AttachmentInfo,
    output_dir: Path,
    skip_conversion: bool = False,
) -> tuple[Path, Path | None]:
    """Save an attachment and optionally convert to text.

    Returns (original_path, converted_path_or_none).
    """
    output_dir.mkdir(parents=True, exist_ok=True)
    original_path = output_dir / attachment.filename

    # Handle duplicate filenames
    counter = 1
    while original_path.exists():
        stem = Path(attachment.filename).stem
        suffix = Path(attachment.filename).suffix
        original_path = output_dir / f"{stem}_{counter}{suffix}"
        counter += 1

    original_path.write_bytes(attachment.data)

    if skip_conversion:
        return original_path, None

    converted_path = _convert(attachment, original_path)
    return original_path, converted_path


def _convert(attachment: AttachmentInfo, original_path: Path) -> Path | None:
    """Attempt to convert an attachment to LLM-readable text."""
    ct = attachment.content_type.lower()
    filename = attachment.filename.lower()

    try:
        if ct == "application/pdf" or filename.endswith(".pdf"):
            return _convert_pdf(attachment.data, original_path)

        if "wordprocessingml" in ct or filename.endswith(".docx"):
            return _convert_docx(attachment.data, original_path)

        if "spreadsheetml" in ct or filename.endswith(".xlsx"):
            return _convert_xlsx(attachment.data, original_path)

        if "presentationml" in ct or filename.endswith(".pptx"):
            return _convert_pptx(attachment.data, original_path)

        if ct == "text/html" or filename.endswith(".html") or filename.endswith(".htm"):
            return _convert_html(attachment.data, original_path)

        if ct == "text/csv" or filename.endswith(".csv"):
            # CSV is already text-friendly, no conversion needed
            return None

        if ct.startswith("text/") or filename.endswith(".txt"):
            # Already text, no conversion needed
            return None

        # Images, archives, etc. — keep as-is
        return None

    except Exception as e:
        console.print(f"[yellow]  Warning: Could not convert {attachment.filename}: {e}[/yellow]")
        return None


def _convert_pdf(data: bytes, original_path: Path) -> Path | None:
    """Convert PDF to text using pdfplumber."""
    import pdfplumber

    text_parts = []
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            page_text = page.extract_text()
            if page_text:
                text_parts.append(f"--- Page {i} ---\n{page_text}")

    if not text_parts:
        return None

    out_path = original_path.with_suffix(original_path.suffix + ".txt")
    out_path.write_text("\n\n".join(text_parts), encoding="utf-8")
    return out_path


def _convert_docx(data: bytes, original_path: Path) -> Path | None:
    """Convert DOCX to text using python-docx."""
    import docx

    doc = docx.Document(io.BytesIO(data))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]

    if not paragraphs:
        return None

    out_path = original_path.with_suffix(original_path.suffix + ".txt")
    out_path.write_text("\n\n".join(paragraphs), encoding="utf-8")
    return out_path


def _convert_xlsx(data: bytes, original_path: Path) -> Path | None:
    """Convert XLSX to CSV using openpyxl."""
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    all_text = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        output = io.StringIO()
        writer = csv.writer(output)

        header = f"=== Sheet: {sheet_name} ==="
        has_data = False

        for row in ws.iter_rows(values_only=True):
            values = [str(cell) if cell is not None else "" for cell in row]
            if any(v.strip() for v in values):
                if not has_data:
                    all_text.append(header)
                    has_data = True
                writer.writerow(values)

        if has_data:
            all_text.append(output.getvalue())

    wb.close()

    if not all_text:
        return None

    out_path = original_path.with_suffix(original_path.suffix + ".csv")
    out_path.write_text("\n".join(all_text), encoding="utf-8")
    return out_path


def _convert_pptx(data: bytes, original_path: Path) -> Path | None:
    """Convert PPTX to text using python-pptx."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    text_parts = []

    for i, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_texts.append(text)
        if slide_texts:
            text_parts.append(f"--- Slide {i} ---\n" + "\n".join(slide_texts))

    if not text_parts:
        return None

    out_path = original_path.with_suffix(original_path.suffix + ".txt")
    out_path.write_text("\n\n".join(text_parts), encoding="utf-8")
    return out_path


def _convert_html(data: bytes, original_path: Path) -> Path | None:
    """Convert HTML attachment to markdown."""
    from bs4 import BeautifulSoup
    from markdownify import markdownify

    html = data.decode("utf-8", errors="replace")
    soup = BeautifulSoup(html, "html.parser")

    for tag in soup.find_all(["style", "script"]):
        tag.decompose()

    md = markdownify(str(soup), heading_style="ATX")
    if not md.strip():
        return None

    out_path = original_path.with_suffix(original_path.suffix + ".md")
    out_path.write_text(md, encoding="utf-8")
    return out_path
